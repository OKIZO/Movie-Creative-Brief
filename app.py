import streamlit as st
import json
from pptx import Presentation
from io import BytesIO

# --- 認証機能 ---
def check_password():
    if "password_correct" not in st.session_state:
        st.session_state["password_correct"] = False
    if st.session_state["password_correct"]:
        return True

    st.write("### 認証が必要です")
    pwd = st.text_input("Password", type="password")
    if st.button("Login"):
        if pwd == st.secrets["password"]:
            st.session_state["password_correct"] = True
            st.rerun()
        else:
            st.error("パスワードが違います")
    return False

# --- テキスト置換用サブ関数 ---
def replace_text_in_text_frame(text_frame, replacements):
    """テキストフレーム内の段落を巡回し、置換を実行する"""
    for paragraph in text_frame.paragraphs:
        if any(key in paragraph.text for key in replacements.keys()):
            full_text = paragraph.text
            for key, val in replacements.items():
                str_val = str(val) if val is not None else ""
                full_text = full_text.replace(key, str_val)
            paragraph.text = full_text

# --- メインのテキスト置換関数 ---
def replace_text_in_presentation(prs, replacements):
    """
    スライド内の全シェイプを巡回し、指定された辞書(replacements)に基づいて
    {{タグ}} を JSON の値に置換します。表（テーブル）の中身も対応。
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            # 1. 通常のテキストボックスの場合
            if shape.has_text_frame:
                replace_text_in_text_frame(shape.text_frame, replacements)
            
            # 2. 表（テーブル）の場合
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        replace_text_in_text_frame(cell.text_frame, replacements)

# --- メイン画面 ---
if check_password():
    st.title("PPT Generator")
    st.caption("HTMLツールからJSONを貼り付けてください")

    json_input = st.text_area("JSON Paste Area", height=250)

    if st.button("パワーポイントを生成"):
        if json_input:
            try:
                data = json.loads(json_input)
                
                # 1. テンプレートの読み込み
                prs = Presentation("template.pptx")
                
                # 2. JSONデータとテンプレートの {{タグ}} を紐づける辞書を作成
                replacements = {
                    "{{映像種別}}": data.get("brief", {}).get("item_name", ""),
                    "{{目的}}": data.get("brief", {}).get("purpose", ""),
                    "{{チャネル}}": data.get("brief", {}).get("channel", ""),
                    "{{対象}}": data.get("brief", {}).get("target", ""),
                    "{{尺}}": data.get("brief", {}).get("duration", ""),
                    
                    "{{Vision}}": data.get("ideology", {}).get("vision", ""),
                    
                    "{{As_is}}": data.get("value", {}).get("behavior_change", {}).get("as_is", ""),
                    "{{To_be}}": data.get("value", {}).get("behavior_change", {}).get("to_be", ""),
                    
                    # リスト形式のものは最初の要素を取得（空の場合は空文字）
                    "{{Benefit}}": data.get("value", {}).get("benefit", [""])[0] if data.get("value", {}).get("benefit") else "",
                    
                    "{{社会背景}}": data.get("context", {}).get("social", {}).get("text", ""),
                    "{{患者インサイト}}": data.get("context", {}).get("patient", {}).get("text", ""),
                    "{{医師インサイト}}": data.get("context", {}).get("doctor", {}).get("text", ""),
                }

                # テーマ案A, B, C の動的マッピング
                axes = data.get("proposed_axes", [])
                if len(axes) > 0:
                    replacements.update({
                        "{{軸案A_軸名}}": axes[0].get("name", ""),
                        "{{軸案A_軸カテゴリ}}": axes[0].get("category", ""),
                        "{{軸案A_主語}}": axes[0].get("subject", ""),
                        "{{軸案A_主役テーマ}}": axes[0].get("theme", ""),
                        "{{軸案A_動画方向性}}": axes[0].get("direction", ""),
                        "{{軸案A_適した型}}": axes[0].get("type", ""),
                    })
                if len(axes) > 1:
                    replacements.update({
                        "{{軸案B_軸名}}": axes[1].get("name", ""),
                        "{{軸案B_軸カテゴリ}}": axes[1].get("category", ""),
                        "{{軸案B_主語}}": axes[1].get("subject", ""),
                        "{{軸案B_主役テーマ}}": axes[1].get("theme", ""),
                        "{{軸案B_動画方向性}}": axes[1].get("direction", ""),
                        "{{軸案B_適した型}}": axes[1].get("type", ""),
                    })
                if len(axes) > 2:
                    replacements.update({
                        "{{軸案C_軸名}}": axes[2].get("name", ""),
                        "{{軸案C_軸カテゴリ}}": axes[2].get("category", ""),
                        "{{軸案C_主語}}": axes[2].get("subject", ""),
                        "{{軸案C_主役テーマ}}": axes[2].get("theme", ""),
                        "{{軸案C_動画方向性}}": axes[2].get("direction", ""),
                        "{{軸案C_適した型}}": axes[2].get("type", ""),
                    })

                # 3. テキストの置換処理を実行
                replace_text_in_presentation(prs, replacements)
                
                # 4. メモリ上に保存してダウンロード可能にする
                ppt_out = BytesIO()
                prs.save(ppt_out)
                ppt_out.seek(0)

                st.success("スライドの生成が完了しました！")
                st.download_button(
                    label="📥 PPTをダウンロード",
                    data=ppt_out,
                    file_name="generated_brief.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"エラーが発生しました: {e}")
